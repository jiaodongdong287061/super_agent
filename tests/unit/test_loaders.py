import pytest
from pathlib import Path
from unittest.mock import patch, MagicMock
from super_agent.knowledge.loaders import get_loader

FIXTURES = Path(__file__).parent.parent.parent / "data" / "raw_docs"


class TestGetLoader:
    def test_pdf_extension(self):
        loader = get_loader(".pdf")
        assert loader is not None
        assert ".pdf" in loader.supported_extensions()

    def test_docx_extension(self):
        loader = get_loader(".docx")
        assert loader is not None

    def test_doc_extension(self):
        loader = get_loader(".doc")
        assert loader is not None

    def test_md_extension(self):
        loader = get_loader(".md")
        assert loader is not None

    def test_html_extension(self):
        loader = get_loader(".html")
        assert loader is not None

    def test_json_extension(self):
        loader = get_loader(".json")
        assert loader is not None

    def test_csv_extension(self):
        loader = get_loader(".csv")
        assert loader is not None

    def test_txt_extension(self):
        loader = get_loader(".txt")
        assert loader is not None

    def test_unknown_extension_raises(self):
        with pytest.raises(ValueError, match="Unsupported"):
            get_loader(".xyz")


class TestMarkdownLoader:
    def test_load_sample(self):
        loader = get_loader(".md")
        docs = loader.load(str(FIXTURES / "sample.md"))
        assert len(docs) > 0
        assert any("MySQL" in d.page_content for d in docs)


class TestPDFLoader:
    def test_is_scanned_page_blank(self):
        from super_agent.knowledge.loaders.pdf import PDFLoader

        loader = PDFLoader()
        mock_page = MagicMock()
        mock_page.rect.width = 595
        mock_page.rect.height = 842
        # page_area = 595 * 842 = 500990, threshold = 500990 * 0.5 / 1000 = 250.5
        assert loader._is_scanned_page("", mock_page) is True
        assert loader._is_scanned_page("   ", mock_page) is True

    def test_is_scanned_page_normal(self):
        from super_agent.knowledge.loaders.pdf import PDFLoader

        loader = PDFLoader()
        mock_page = MagicMock()
        mock_page.rect.width = 595
        mock_page.rect.height = 842
        long_text = "x" * 300
        assert loader._is_scanned_page(long_text, mock_page) is False

    def test_is_scanned_page_ocr_disabled(self):
        from super_agent.knowledge.loaders.pdf import PDFLoader
        from super_agent.config import OCRConfig

        loader = PDFLoader()
        mock_page = MagicMock()
        mock_page.rect.width = 595
        mock_page.rect.height = 842
        with patch("super_agent.knowledge.loaders.pdf.settings") as mock_settings:
            mock_settings.ocr = OCRConfig(enabled=False)
            assert loader._is_scanned_page("", mock_page) is False

    def test_ocr_engine_not_installed(self):
        import super_agent.knowledge.loaders.pdf as pdf_module

        pdf_module._PADDLEOCR_AVAILABLE = None
        with patch.dict("sys.modules", {"paddleocr": None}):
            assert pdf_module._check_paddleocr() is False
            assert pdf_module._get_ocr_engine() is None
        pdf_module._PADDLEOCR_AVAILABLE = None
        pdf_module._get_ocr_engine.cache_clear()

    def test_ocr_page_returns_empty_when_no_engine(self):
        from super_agent.knowledge.loaders.pdf import PDFLoader

        loader = PDFLoader()
        mock_page = MagicMock()
        with patch("super_agent.knowledge.loaders.pdf._get_ocr_engine", return_value=None):
            assert loader._ocr_page(mock_page) == ""


class TestExcelLoaderNoTopicTags:
    def test_metadata_no_topic_tags_key(self, tmp_path):
        """ExcelLoader 不应在 metadata 中硬编码 topic_tags"""
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["名称", "值"])
        ws.append(["A", "1"])
        ws.append(["B", "2"])
        f = tmp_path / "test.xlsx"
        wb.save(str(f))

        from super_agent.knowledge.loaders.excel import ExcelLoader

        loader = ExcelLoader()
        docs = loader.load(str(f))
        assert len(docs) > 0
        assert "topic_tags" not in docs[0].metadata


class TestPPTLoader:
    def test_pptx_extension_registered(self):
        loader = get_loader(".pptx")
        assert loader is not None
        assert ".pptx" in loader.supported_extensions()

    def test_ppt_extension_registered(self):
        loader = get_loader(".ppt")
        assert loader is not None
        assert ".ppt" in loader.supported_extensions()

    def test_load_pptx_sample(self, tmp_path):
        """用 python-pptx 构造一个最小 pptx 文件并加载"""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "测试标题"
        body = slide.placeholders[1]
        body.text = "测试正文内容"

        f = tmp_path / "test.pptx"
        prs.save(str(f))

        from super_agent.knowledge.loaders.ppt import PPTLoader

        loader = PPTLoader()
        docs = loader.load(str(f))
        assert len(docs) == 1
        assert "测试标题" in docs[0].page_content
        assert "测试正文内容" in docs[0].page_content

    def test_pptx_metadata(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "标题"

        f = tmp_path / "test.pptx"
        prs.save(str(f))

        from super_agent.knowledge.loaders.ppt import PPTLoader

        loader = PPTLoader()
        docs = loader.load(str(f))
        assert len(docs) == 1
        meta = docs[0].metadata
        assert meta["source"] == str(f)
        assert meta["slide_number"] == 1
        assert meta["total_slides"] == 1
        assert "has_notes" in meta
        assert "has_tables" in meta
        assert "has_images" in meta
        assert "ocr_used" in meta

    def test_multiple_slides(self, tmp_path):
        """多张幻灯片各产出 1 个 Document"""
        from pptx import Presentation

        prs = Presentation()
        for i in range(3):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"第{i + 1}页"
            body = slide.placeholders[1]
            body.text = f"内容{i + 1}"

        f = tmp_path / "multi.pptx"
        prs.save(str(f))

        from super_agent.knowledge.loaders.ppt import PPTLoader

        loader = PPTLoader()
        docs = loader.load(str(f))
        assert len(docs) == 3
        assert docs[0].metadata["slide_number"] == 1
        assert docs[1].metadata["slide_number"] == 2
        assert docs[2].metadata["slide_number"] == 3
        assert docs[0].metadata["total_slides"] == 3

    def test_table_extraction(self, tmp_path):
        """表格提取为 Markdown 格式"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        rows, cols = 2, 2
        table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "名称"
        table.cell(0, 1).text = "值"
        table.cell(1, 0).text = "CPU"
        table.cell(1, 1).text = "90%"

        f = tmp_path / "table.pptx"
        prs.save(str(f))

        from super_agent.knowledge.loaders.ppt import PPTLoader

        loader = PPTLoader()
        docs = loader.load(str(f))
        assert len(docs) == 1
        assert "|" in docs[0].page_content
        assert "CPU" in docs[0].page_content
        assert docs[0].metadata["has_tables"] is True

    def test_notes_extraction(self, tmp_path):
        """演讲者备注以 [备注] 前缀提取"""
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "标题"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "这是备注内容"

        f = tmp_path / "notes.pptx"
        prs.save(str(f))

        from super_agent.knowledge.loaders.ppt import PPTLoader

        loader = PPTLoader()
        docs = loader.load(str(f))
        assert len(docs) == 1
        assert "[备注]" in docs[0].page_content
        assert "这是备注内容" in docs[0].page_content
        assert docs[0].metadata["has_notes"] is True

    def test_unsupported_extension_raises(self):
        from super_agent.knowledge.loaders.ppt import PPTLoader

        loader = PPTLoader()
        with pytest.raises(ValueError, match="Unsupported extension"):
            loader.load("test.txt")
