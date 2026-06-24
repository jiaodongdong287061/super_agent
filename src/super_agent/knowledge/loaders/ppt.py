from __future__ import annotations

import logging
import subprocess
import tempfile
from pathlib import Path

from langchain_core.documents import Document

from super_agent.config import settings
from super_agent.knowledge.loaders.base import BaseLoader

logger = logging.getLogger(__name__)

_PADDLEOCR_AVAILABLE: bool | None = None


def _check_paddleocr() -> bool:
    global _PADDLEOCR_AVAILABLE
    if _PADDLEOCR_AVAILABLE is None:
        try:
            import paddleocr  # noqa: F401

            _PADDLEOCR_AVAILABLE = True
        except ImportError:
            _PADDLEOCR_AVAILABLE = False
            logger.warning(
                "paddleocr is not installed. Image OCR in PPT will be skipped. "
                "Install with: uv sync --extra ml"
            )
    return _PADDLEOCR_AVAILABLE


def _get_ocr_engine():
    from functools import lru_cache

    @lru_cache(maxsize=1)
    def _cached():
        if not _check_paddleocr():
            return None
        from paddleocr import PaddleOCR

        return PaddleOCR(use_gpu=settings.ocr.use_gpu, lang=settings.ocr.lang, show_log=False)

    return _cached()


class PPTLoader(BaseLoader):
    def load(self, source: str) -> list[Document]:
        ext = Path(source).suffix.lower()
        if ext == ".pptx":
            return self._load_pptx(source)
        if ext == ".ppt":
            return self._load_ppt(source)
        raise ValueError(f"Unsupported extension: {ext}")

    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    def _load_ppt(self, source: str) -> list[Document]:
        """Convert .ppt to .pptx via LibreOffice, then load as pptx."""
        with tempfile.TemporaryDirectory() as tmp_dir:
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        tmp_dir,
                        source,
                    ],
                    timeout=60,
                    check=True,
                    capture_output=True,
                )
            except FileNotFoundError:
                raise RuntimeError(
                    "LibreOffice is required for .ppt conversion but not found. "
                    "Install LibreOffice or convert .ppt to .pptx manually."
                )
            except subprocess.TimeoutExpired:
                raise RuntimeError(
                    f"LibreOffice conversion timed out (60s) for {source}"
                )

            pptx_path = Path(tmp_dir) / (Path(source).stem + ".pptx")
            if not pptx_path.exists():
                raise RuntimeError(f"LibreOffice conversion failed: {pptx_path} not found")

            return self._load_pptx(str(pptx_path))

    def _load_pptx(self, source: str) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(source)
        total = len(prs.slides)
        docs: list[Document] = []

        for idx, slide in enumerate(prs.slides, start=1):
            content_parts: list[str] = []
            has_notes = False
            has_tables = False
            has_images = False
            ocr_used = False

            for shape in slide.shapes:
                # 文本
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        content_parts.append(text)

                # 表格
                if shape.has_table:
                    has_tables = True
                    table_text = self._extract_table(shape.table)
                    if table_text:
                        content_parts.append(table_text)

                # 图片 OCR
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    has_images = True
                    if settings.ocr.enabled:
                        ocr_text = self._ocr_shape(shape)
                        if ocr_text:
                            ocr_used = True
                            content_parts.append(f"[OCR] {ocr_text}")

            # 演讲者备注
            notes_text = self._extract_notes(slide)
            if notes_text:
                has_notes = True
                content_parts.append(f"[备注] {notes_text}")

            page_content = "\n\n".join(content_parts)
            if not page_content.strip():
                continue

            docs.append(
                Document(
                    page_content=page_content,
                    metadata={
                        "source": source,
                        "slide_number": idx,
                        "total_slides": total,
                        "has_notes": has_notes,
                        "has_tables": has_tables,
                        "has_images": has_images,
                        "ocr_used": ocr_used,
                    },
                )
            )

        return docs

    @staticmethod
    def _extract_table(table) -> str:
        """将表格转换为 Markdown 格式。"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if not rows:
            return ""
        header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
        rows.insert(1, header_sep)
        return "\n".join(rows)

    @staticmethod
    def _extract_notes(slide) -> str:
        """提取演讲者备注。"""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            text = notes_frame.text.strip()
            return text
        return ""

    def _ocr_shape(self, shape) -> str:
        """对 PPT 中的图片执行 OCR。"""
        engine = _get_ocr_engine()
        if engine is None:
            return ""
        try:
            image = shape.image
            img_bytes = image.blob
            result = engine.ocr(img_bytes, cls=True)
            if not result or not result[0]:
                return ""
            lines = [line[1][0] for line in result[0]]
            return "\n".join(lines)
        except Exception:
            logger.warning("OCR failed for a PPT image, skipping", exc_info=True)
            return ""
